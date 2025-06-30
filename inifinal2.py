import customtkinter as ctk
import threading
import json
import os
import re
from tkinter import filedialog
import io

# --- Pustaka Pihak Ketiga ---
try:
    import requests
    import fitz  # PyMuPDF
    from pptx import Presentation
    from PIL import Image
    import pytesseract
except ImportError as e:
    print(f"Error: Pustaka yang dibutuhkan hilang -> {e.name}")
    print("Silakan jalankan: pip install --upgrade customtkinter requests PyMuPDF python-pptx pytesseract Pillow")
    exit()

# --- Konfigurasi Aplikasi ---
try:
    # Arahkan ke instalasi Tesseract Anda jika perlu
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except Exception:
    print("Peringatan: Tesseract tidak ditemukan. Fitur OCR gambar tidak akan berfungsi.")

try:
    # Menggunakan tema kustom jika ada
    ctk.set_default_color_theme("theme.json")
except FileNotFoundError:
    print("Peringatan: file 'theme.json' tidak ditemukan. Menggunakan tema default.")
    ctk.set_default_color_theme("blue")

ctk.set_appearance_mode("Dark")

# Definisi Font & Warna
APP_FONT_FAMILY = "Arial"
H1_FONT = (APP_FONT_FAMILY, 48, "bold")
H2_FONT = (APP_FONT_FAMILY, 24, "bold")
H3_FONT = (APP_FONT_FAMILY, 20, "bold")
BODY_FONT = (APP_FONT_FAMILY, 14)
BUTTON_FONT = (APP_FONT_FAMILY, 16, "bold")
SMALL_FONT = (APP_FONT_FAMILY, 12)
SUCCESS_COLOR = "#00A896"
ERROR_COLOR = "#D90429"
GOLD_ACCENT_COLOR = ("#C09A3E", "#D4AF37")
AI_BUBBLE_COLOR = ("#E5E5E5", "#3A3A3A")


# --- Model AI Gemini & Fungsi Helper ---
# Pastikan API Key Anda dimasukkan di sini
API_KEY = "AIzaSyAX4tK_mDxeBC4T08WS_1Iz0JSjyPSrGKo" # @param {type:"string"}
API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={API_KEY}"

def get_gemini_response(prompt, content_text=None):
    if not API_KEY:
        return "Fitur AI dinonaktifkan. Silakan masukkan Kunci API Anda di dalam source code."
    headers = {'Content-Type': 'application/json'}
    full_prompt = f"{prompt}\n\nBerikut kontennya:\n{content_text}" if content_text else prompt
    payload = {"contents": [{"parts": [{"text": full_prompt}]}]}
    try:
        response = requests.post(API_URL, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        result = response.json()
        if 'candidates' in result and result['candidates']:
            # Menghapus markdown dari respons agar bersih saat ditampilkan di Label
            text = result['candidates'][0]['content']['parts'][0]['text']
            return re.sub(r'(\*\*|\*)', '', text)
        else:
            return "Error: Respons AI tidak valid atau kosong."
    except requests.exceptions.RequestException as e:
        return f"Error jaringan saat menghubungi AI: {e}"
    except Exception as e:
        return f"Error tak terduga saat menghubungi AI: {e}"

# --- Komponen-komponen UI ---
class QuizWindow(ctk.CTkToplevel):
    def __init__(self, parent, quiz_data):
        super().__init__(parent)
        self.title("Grind On - Kuis")
        self.geometry("800x650")
        self.transient(parent)
        self.grab_set()

        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)
        self.grid_rowconfigure(2, weight=1)
        
        self.quiz_data = quiz_data
        self.current_question_index = 0
        self.score = 0
        self.selected_option_index = None
        
        self._create_widgets()
        self.display_question()
        self.bind("<Configure>", self.update_wraplength)

    def update_wraplength(self, event=None):
        new_width = self.question_frame.winfo_width() - 50 
        if new_width > 1: self.question_label.configure(wraplength=new_width)

    def _create_widgets(self):
        self.top_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.top_frame.grid(row=0, column=0, sticky="ew", padx=20, pady=10)
        self.question_number_label = ctk.CTkLabel(self.top_frame, font=H3_FONT)
        self.question_number_label.pack(side="left")
        self.score_label = ctk.CTkLabel(self.top_frame, text="Skor: 0", font=H3_FONT)
        self.score_label.pack(side="right")

        self.question_frame = ctk.CTkFrame(self, corner_radius=10)
        self.question_frame.grid(row=1, column=0, sticky="nsew", padx=20, pady=10)
        self.question_frame.grid_columnconfigure(0, weight=1)
        self.question_frame.grid_rowconfigure(0, weight=1)
        self.question_label = ctk.CTkLabel(self.question_frame, font=H2_FONT, text="")
        self.question_label.grid(row=0, column=0, padx=25, pady=25, sticky="nsew")

        self.options_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.options_frame.grid(row=2, column=0, sticky="nsew", padx=20, pady=10)
        self.options_frame.grid_columnconfigure((0, 1), weight=1)
        self.options_frame.grid_rowconfigure((0, 1), weight=1)
        
        self.option_buttons = []
        for i in range(4):
            # --- PERBAIKAN: Menghapus argumen 'wraplength' dan 'justify' ---
            button = ctk.CTkButton(self.options_frame, font=BUTTON_FONT, height=70, corner_radius=8, command=lambda i=i: self.select_option(i))
            button.grid(row=i//2, column=i%2, padx=5, pady=5, sticky="nsew")
            self.option_buttons.append(button)

        self.bottom_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.bottom_frame.grid(row=3, column=0, sticky="ew", padx=20, pady=(10, 20))
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.submit_button = ctk.CTkButton(self.bottom_frame, text="Submit Jawaban", font=BUTTON_FONT, height=45, command=self.submit_answer, state="disabled")
        self.submit_button.grid(row=0, column=0, sticky="ew")
        
        self.next_button = ctk.CTkButton(self.bottom_frame, text="Pertanyaan Berikutnya", font=BUTTON_FONT, height=45, command=self.next_question)
        
        self.explanation_frame = ctk.CTkFrame(self, corner_radius=10)
        self.explanation_label = ctk.CTkLabel(self.explanation_frame, text="", justify="left", font=BODY_FONT)
        self.explanation_label.pack(padx=15, pady=15, fill="x", expand=True)

    def display_question(self):
        self.explanation_frame.grid_forget()
        self.next_button.grid_forget()
        self.submit_button.grid(row=0, column=0, sticky="ew")
        self.submit_button.configure(state="disabled")
        self.selected_option_index = None

        q_data = self.quiz_data[self.current_question_index]
        self.question_number_label.configure(text=f"Soal {self.current_question_index + 1}/{len(self.quiz_data)}")
        self.question_label.configure(text=q_data.get("question", "Pertanyaan tidak tersedia"))
        
        self.update_wraplength()
        
        options = q_data.get("options", [])
        for i, button in enumerate(self.option_buttons):
            option_text = options[i] if i < len(options) else f"Opsi {i+1}"
            button.configure(text=option_text, state="normal", fg_color=ctk.ThemeManager.theme["CTkButton"]["fg_color"])

    def select_option(self, index):
        self.selected_option_index = index
        self.submit_button.configure(state="normal")
        for i, button in enumerate(self.option_buttons):
            is_selected = (i == index)
            button.configure(fg_color=ctk.ThemeManager.theme["CTkButton"]["hover_color"] if is_selected else ctk.ThemeManager.theme["CTkButton"]["fg_color"])

    def submit_answer(self):
        q_data = self.quiz_data[self.current_question_index]
        correct_index = q_data.get("correct_answer_index")
        
        for button in self.option_buttons: button.configure(state="disabled")

        if self.selected_option_index == correct_index:
            self.score += 1
            self.option_buttons[self.selected_option_index].configure(fg_color=SUCCESS_COLOR)
        else:
            self.option_buttons[self.selected_option_index].configure(fg_color=ERROR_COLOR)
            if correct_index is not None and 0 <= correct_index < 4:
                self.option_buttons[correct_index].configure(fg_color=SUCCESS_COLOR)
        
        self.score_label.configure(text=f"Skor: {self.score}")
        
        explanation_text = f"Penjelasan:\n{q_data.get('explanation', 'Penjelasan tidak tersedia.')}"
        self.explanation_label.configure(text=explanation_text)
        self.explanation_label.configure(wraplength=self.explanation_frame.winfo_width() - 30)

        self.explanation_frame.grid(row=4, column=0, padx=20, pady=10, sticky="ew")
        
        self.submit_button.grid_forget()
        self.next_button.grid(row=0, column=0, sticky="ew")

    def next_question(self):
        self.current_question_index += 1
        if self.current_question_index < len(self.quiz_data):
            self.display_question()
        else:
            self.show_results()

    def show_results(self):
        for widget in self.winfo_children(): widget.destroy()
        self.grid_rowconfigure((0,1,2), weight=1, uniform="results")
        self.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(self, text="Kuis Selesai!", font=(APP_FONT_FAMILY, 40, "bold")).grid(row=0, column=0)
        ctk.CTkLabel(self, text=f"Skor Akhir: {self.score} / {len(self.quiz_data)}", font=H2_FONT).grid(row=1, column=0)
        ctk.CTkButton(self, text="Tutup", font=BUTTON_FONT, command=self.destroy, height=45).grid(row=2, column=0, padx=100, pady=20, sticky="ew")


class SmartStudyApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("IHTLI-Study App")
        self.geometry("1100x720")
        self.minsize(900, 600)

        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self._create_sidebar()
        self._create_main_frame()
        self._create_popup_frames()
        self.select_feature("pomodoro")

    def _create_sidebar(self):
        self.sidebar_frame = ctk.CTkFrame(self, width=220, corner_radius=0)
        self.sidebar_frame.grid(row=0, column=0, rowspan=4, sticky="nsew")
        self.sidebar_frame.grid_rowconfigure(6, weight=1)

        self.logo_label = ctk.CTkLabel(self.sidebar_frame, text="IHTLI-Study App", font=H2_FONT)
        self.logo_label.grid(row=0, column=0, padx=20, pady=(30, 20))

        buttons_info = {
            "pomodoro": "Pomodoro", "chatbot": "Ask Me AI", "summarizer": "Summarizer",
            "explainer": "Explainer", "quiz": "Grind On"
        }
        self.sidebar_buttons = {}
        for i, (name, text) in enumerate(buttons_info.items(), start=1):
            button = ctk.CTkButton(self.sidebar_frame, text=text, font=BUTTON_FONT, anchor="w", 
                                   height=45, corner_radius=8, command=lambda n=name: self.select_feature(n))
            button.grid(row=i, column=0, padx=20, pady=8, sticky="ew")
            self.sidebar_buttons[name] = button
            
        self.theme_switch = ctk.CTkSwitch(self.sidebar_frame, text="Light Mode", font=BODY_FONT, command=self.toggle_theme)
        self.theme_switch.grid(row=8, column=0, padx=20, pady=(20, 20), sticky="w")
        
    def toggle_theme(self):
        if self.theme_switch.get() == 1:
            ctk.set_appearance_mode("Light")
            self.theme_switch.configure(text="Dark Mode")
        else:
            ctk.set_appearance_mode("Dark")
            self.theme_switch.configure(text="Light Mode")
        
        if "pomodoro" in self.popup_frames and self.popup_frames["pomodoro"].winfo_exists():
            self.popup_frames["pomodoro"].on_theme_change()

    def _create_main_frame(self):
        self.main_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.main_frame.grid(row=0, column=1, padx=20, pady=20, sticky="nsew")
        self.main_frame.grid_columnconfigure(0, weight=1)
        self.main_frame.grid_rowconfigure(1, weight=1)

        self.feature_title_frame = ctk.CTkFrame(self.main_frame, corner_radius=10, height=60)
        self.feature_title_frame.grid(row=0, column=0, pady=(0, 10), sticky="ew")
        
        self.feature_title_label = ctk.CTkLabel(self.feature_title_frame, text="", font=H2_FONT)
        self.feature_title_label.pack(side="left", padx=20, pady=10)

    def _create_popup_frames(self):
        self.popup_frames = {
            "pomodoro": PomodoroFrame(self.main_frame),
            "chatbot": ChatbotFrame(self.main_frame),
            "summarizer": DocumentFeatureFrame(self, "summarize"),
            "explainer": DocumentFeatureFrame(self, "explain"),
            "quiz": DocumentFeatureFrame(self, "quiz")
        }
        for frame in self.popup_frames.values():
            frame.grid(row=1, column=0, sticky="nsew", in_=self.main_frame)
            frame.grid_remove() 

    def select_feature(self, name):
        titles = {
            "pomodoro": "Pomodoro Timer", "chatbot": "Ask Me AI Assistant", "summarizer": "Document Summarizer",
            "explainer": "Concept Explainer", "quiz": "Quiz Generator (Grind On)"
        }
        self.feature_title_label.configure(text=titles.get(name, "Fitur"))

        for button_name, button in self.sidebar_buttons.items():
             button.configure(fg_color=ctk.ThemeManager.theme["CTkButton"]["hover_color"] if button_name == name else "transparent")

        for frame_name, frame in self.popup_frames.items():
            if frame_name == name:
                frame.grid()
                frame.lift()
            else:
                frame.grid_remove()
                
class ChatbotFrame(ctk.CTkFrame):
    def __init__(self, parent):
        super().__init__(parent, fg_color="transparent")
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self.chat_history_frame = ctk.CTkScrollableFrame(self, corner_radius=10)
        self.chat_history_frame.grid(row=0, column=0, sticky="nsew", pady=(0,10))
        self.chat_history_frame.grid_columnconfigure(0, weight=1)
        self.chat_history_frame.bind("<Configure>", self._on_resize)
        self.last_width = self.chat_history_frame.winfo_width()


        self.input_frame = ctk.CTkFrame(self, corner_radius=10)
        self.input_frame.grid(row=1, column=0, sticky="ew")
        self.input_frame.grid_columnconfigure(0, weight=1)

        self.user_input = ctk.CTkEntry(self.input_frame, placeholder_text="Ketik pesan Anda...", font=BODY_FONT, border_width=0, fg_color="transparent")
        self.user_input.grid(row=0, column=0, padx=15, pady=5, sticky="ew")
        self.user_input.bind("<Return>", lambda event: self.send_message())

        self.send_button = ctk.CTkButton(self.input_frame, text="Kirim", font=BUTTON_FONT, width=80, height=35, command=self.send_message)
        self.send_button.grid(row=0, column=1, padx=(0, 5), pady=5)
        
        self._add_message_bubble("bot", "Halo! Saya adalah asisten AI. Apa yang bisa saya bantu?")

    def _on_resize(self, event=None):
        """Memperbarui wraplength semua label saat jendela diubah ukurannya."""
        if abs(self.chat_history_frame.winfo_width() - self.last_width) < 10:
             return 
        
        self.last_width = self.chat_history_frame.winfo_width()
        for container in self.chat_history_frame.winfo_children():
            if len(container.winfo_children()) > 1 and isinstance(container.winfo_children()[1], ctk.CTkFrame):
                bubble_frame = container.winfo_children()[1]
                if len(bubble_frame.winfo_children()) > 0:
                    content_label = bubble_frame.winfo_children()[0]
                    if isinstance(content_label, ctk.CTkLabel):
                        new_wraplength = self.chat_history_frame.winfo_width() - 100
                        if new_wraplength > 0:
                            content_label.configure(wraplength=new_wraplength)
    
    def _create_bubble(self, sender, message):
        """Membuat gelembung pesan yang dinamis dengan label pengirim."""
        is_user = (sender == "user")
        anchor = "e" if is_user else "w"
        
        container = ctk.CTkFrame(self.chat_history_frame, fg_color="transparent")
        container.pack(fill="x", padx=5, pady=(8,0))

        sender_text = "Anda" if is_user else "AI"
        sender_label = ctk.CTkLabel(container, text=sender_text, font=(SMALL_FONT[0], SMALL_FONT[1], "bold"))
        sender_label.pack(anchor=anchor, padx=10, pady=(0,2))
        
        bubble_frame = ctk.CTkFrame(container, 
                                    fg_color=GOLD_ACCENT_COLOR if is_user else AI_BUBBLE_COLOR,
                                    corner_radius=15)
        bubble_frame.pack(anchor=anchor, padx=(50, 5) if is_user else (5, 50))

        initial_wraplength = max(1, self.chat_history_frame.winfo_width() - 100)
        content_label = ctk.CTkLabel(bubble_frame, text=message, font=BODY_FONT, justify="left",
                                     wraplength=initial_wraplength)
        content_label.pack(padx=12, pady=10)
        
        if is_user:
            content_label.configure(text_color=("#1A1A1A", "#1A1A1A"))

    def _add_message_bubble(self, sender, message):
        if hasattr(self, 'typing_bubble_frame'):
            self.typing_bubble_frame.destroy()
            del self.typing_bubble_frame

        self._create_bubble(sender, message)
        
        self.after(100, self.chat_history_frame._parent_canvas.yview_moveto, 1.0)

    def _create_typing_indicator(self):
        self.typing_bubble_frame = ctk.CTkFrame(self.chat_history_frame, fg_color="transparent")
        bubble = ctk.CTkLabel(self.typing_bubble_frame, text="AI sedang mengetik...", font=(APP_FONT_FAMILY, 14, "italic"),
                              fg_color=AI_BUBBLE_COLOR, corner_radius=15,
                              text_color=("gray50", "gray70"))
        bubble.pack(anchor="w", padx=5, pady=5, ipady=5, ipadx=10)
        self.typing_bubble_frame.pack(fill="x", padx=5, pady=4)
        self.after(100, self.chat_history_frame._parent_canvas.yview_moveto, 1.0)

    def send_message(self):
        user_message = self.user_input.get()
        if not user_message.strip(): return
        
        self._add_message_bubble("user", user_message)
        self.user_input.delete(0, "end")
        self._create_typing_indicator()
        
        threading.Thread(target=self._get_ai_response, args=(user_message,), daemon=True).start()

    def _get_ai_response(self, message):
        prompt = "Anda adalah asisten yang ramah dan membantu. Jawab pertanyaan dengan jelas dan ringkas. Jangan gunakan markdown."
        response = get_gemini_response(prompt, message)
        self.after(0, lambda: self._add_message_bubble("bot", response))

class DocumentFeatureFrame(ctk.CTkFrame):
    def __init__(self, parent_app, feature_type):
        super().__init__(parent_app.main_frame, fg_color="transparent")
        self.parent_app = parent_app
        self.feature_type = feature_type
        self.file_content = None

        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)
        self._create_widgets()

    def _create_widgets(self):
        top_frame = ctk.CTkFrame(self, corner_radius=10)
        top_frame.grid(row=0, column=0, pady=(0, 10), sticky="ew")
        top_frame.grid_columnconfigure(1, weight=1)

        self.upload_button = ctk.CTkButton(top_frame, text="Upload File", font=BUTTON_FONT, height=40, command=self.upload_file)
        self.upload_button.grid(row=0, column=0, padx=10, pady=10)

        self.file_label = ctk.CTkLabel(top_frame, text="Belum ada file yang di-upload", text_color="gray", font=SMALL_FONT, anchor="w")
        self.file_label.grid(row=0, column=1, padx=10, sticky="ew")

        button_text = "Mulai Kuis" if self.feature_type == 'quiz' else "Proses Dokumen"
        self.run_button = ctk.CTkButton(top_frame, text=button_text, font=BUTTON_FONT, height=40, command=self.run_feature, state="disabled")
        self.run_button.grid(row=0, column=2, padx=10, pady=10)

        self.output_scroll_frame = ctk.CTkScrollableFrame(self, corner_radius=10)
        self.output_scroll_frame.grid(row=1, column=0, sticky="nsew")
        self.output_scroll_frame.grid_columnconfigure(0, weight=1)
        
        self.output_scroll_frame.bind("<Configure>", self._on_resize)
        self.last_width = self.output_scroll_frame.winfo_width()
        
        self._display_message("Sistem", "Upload dokumen untuk memulai analisis AI.")

    def _on_resize(self, event=None):
        if abs(self.output_scroll_frame.winfo_width() - self.last_width) < 10:
             return
        
        self.last_width = self.output_scroll_frame.winfo_width()
        for container in self.output_scroll_frame.winfo_children():
            if len(container.winfo_children()) > 1 and isinstance(container.winfo_children()[1], ctk.CTkFrame):
                bubble_frame = container.winfo_children()[1]
                if len(bubble_frame.winfo_children()) > 0:
                    content_label = bubble_frame.winfo_children()[0]
                    if isinstance(content_label, ctk.CTkLabel):
                        new_wraplength = self.output_scroll_frame.winfo_width() - 40
                        if new_wraplength > 0:
                            content_label.configure(wraplength=new_wraplength)

    def _display_message(self, sender, message):
        """Fungsi terpusat untuk menampilkan pesan dalam format bubble."""
        for widget in self.output_scroll_frame.winfo_children():
            widget.destroy()

        container = ctk.CTkFrame(self.output_scroll_frame, fg_color="transparent")
        container.pack(fill="x", padx=5, pady=(8,0))
        
        sender_label = ctk.CTkLabel(container, text=sender, font=(SMALL_FONT[0], SMALL_FONT[1], "bold"))
        sender_label.pack(anchor="w", padx=10, pady=(0,2))
        
        bubble_frame = ctk.CTkFrame(container, fg_color=AI_BUBBLE_COLOR, corner_radius=15)
        bubble_frame.pack(anchor="w", padx=5, fill="x")
        
        initial_wraplength = max(1, self.output_scroll_frame.winfo_width() - 40)
        content_label = ctk.CTkLabel(bubble_frame, text=message, font=BODY_FONT, justify="left",
                                     wraplength=initial_wraplength)
        content_label.pack(padx=12, pady=10, fill="x")

    def upload_file(self):
        filetypes = (("Dokumen & Gambar", "*.pdf *.pptx *.png *.jpg *.jpeg *.bmp *.tiff"),)
        filepath = filedialog.askopenfilename(title="Pilih File", filetypes=filetypes)
        if not filepath: return
        
        filename = os.path.basename(filepath)
        self.file_label.configure(text=filename, text_color=ctk.ThemeManager.theme["CTkLabel"]["text_color"])
        self._display_message("Sistem", f"Menganalisis file '{filename}'...")
        self.run_button.configure(state="disabled")
        
        threading.Thread(target=self._extract_text, args=(filepath,), daemon=True).start()

    def _extract_text(self, file_path):
        try:
            ext = os.path.splitext(file_path)[1].lower()
            text = ""
            if ext == '.pdf':
                with fitz.open(file_path) as doc:
                    if doc.is_encrypted: raise ValueError("PDF ini dilindungi kata sandi.")
                    for page in doc:
                        text += page.get_text("text", sort=True) + "\n"
                        for img_info in page.get_images(full=True):
                            try:
                                base_image = doc.extract_image(img_info[0])
                                image_bytes = base_image["image"]
                                image = Image.open(io.BytesIO(image_bytes))
                                ocr_text = pytesseract.image_to_string(image)
                                if ocr_text.strip():
                                    text += f"\n--- Teks dari Gambar (OCR) ---\n{ocr_text}\n--- Akhir Teks Gambar ---\n"
                            except Exception:
                                continue
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                text = pytesseract.image_to_string(Image.open(file_path))
            else:
                raise ValueError(f"Format file tidak didukung: {ext}")
            
            cleaned_text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
            if not cleaned_text or len(cleaned_text) < 20:
                raise ValueError("Tidak dapat menemukan teks yang signifikan dalam file.")
            
            self.file_content = cleaned_text
            self.after(0, self._on_extraction_complete, os.path.basename(file_path))
        except Exception as e:
            self.after(0, self._display_message, "Error", f"Gagal memproses file:\n\n{e}")

    def _on_extraction_complete(self, filename):
        action = "memulai kuis" if self.feature_type == 'quiz' else "memprosesnya"
        self._display_message("Sistem", f"Sukses! File '{filename}' telah dianalisis.\n\nAnda sekarang bisa {action}.")
        self.run_button.configure(state="normal")

    def run_feature(self):
        if not self.file_content:
            self._display_message("Error", "Tidak ada konten file untuk diproses.")
            return
            
        self._display_message("AI", "Sedang memproses... Mohon tunggu.")
        self.run_button.configure(state="disabled")
        threading.Thread(target=self._get_ai_result, daemon=True).start()

    def _get_ai_result(self):
        prompts = {
            "summarize": "Anda adalah seorang analis ahli. Buat rangkuman yang mendalam dan terperinci dari teks berikut. Jangan gunakan markdown.",
            "explain": "Anda adalah seorang dosen ahli. Identifikasi semua konsep penting dalam teks berikut dan berikan penjelasan yang komprehensif untuk setiap konsep. Jangan gunakan markdown.",
            "quiz": "Anda adalah pembuat kuis. Buat 10 pertanyaan pilihan ganda dari teks ini. Respons HARUS HANYA dalam format JSON array yang valid. Setiap objek dalam array harus memiliki kunci: 'question' (string), 'options' (array berisi 4 string), 'correct_answer_index' (integer dari 0-3), dan 'explanation' (string singkat)."
        }
        prompt = prompts.get(self.feature_type)
        response_text = get_gemini_response(prompt, self.file_content)

        if self.feature_type == 'quiz':
            try:
                match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
                if match:
                    json_str = match.group(1)
                else:
                    json_str = response_text
                
                quiz_data = json.loads(json_str)
                if not isinstance(quiz_data, list) or len(quiz_data) < 1:
                    raise ValueError("Struktur JSON tidak sesuai (bukan list atau kosong).")
                
                self.after(0, lambda: QuizWindow(self.parent_app, quiz_data))
                self.after(0, self._display_message, "Sistem", "Sukses! Kuis telah dimulai di jendela baru.")
            except (json.JSONDecodeError, ValueError) as e:
                error_msg = f"Gagal membuat kuis.\n\nAI memberikan respons dalam format yang tidak valid.\n\nDetail Error:\n{e}"
                self.after(0, self._display_message, "Error", error_msg)
            finally:
                self.after(0, self.run_button.configure, {"state": "normal"})
        else:
            self.after(0, lambda: (
                self._display_message("AI", response_text), 
                self.run_button.configure(state="normal")
            ))

class PomodoroFrame(ctk.CTkFrame):
    def __init__(self, parent):
        super().__init__(parent, fg_color="transparent")
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)

        self.timer_running = False
        self.timer_id = None
        self.current_mode = "Belajar"
        
        self.study_duration_min = ctk.StringVar(value="25")
        self.break_duration_min = ctk.StringVar(value="5")
        
        self.remaining_time = int(self.study_duration_min.get()) * 60
        
        self._create_widgets()
        self.on_theme_change()

    def _create_widgets(self):
        settings_frame = ctk.CTkFrame(self, fg_color="transparent")
        settings_frame.grid(row=0, column=0, pady=(0, 20))
        ctk.CTkLabel(settings_frame, text="Waktu Belajar (menit):", font=BODY_FONT).pack(side="left", padx=(0, 5))
        ctk.CTkOptionMenu(settings_frame, variable=self.study_duration_min, font=BODY_FONT, values=[str(i) for i in [5, 15, 25, 30, 45, 60]], command=self.reset_timer).pack(side="left", padx=(0, 20))
        ctk.CTkLabel(settings_frame, text="Waktu Istirahat (menit):", font=BODY_FONT).pack(side="left", padx=(0, 5))
        ctk.CTkOptionMenu(settings_frame, variable=self.break_duration_min, font=BODY_FONT, values=[str(i) for i in [1, 5, 10, 15, 20]], command=self.reset_timer).pack(side="left")

        timer_frame = ctk.CTkFrame(self, fg_color="transparent")
        timer_frame.grid(row=1, column=0, sticky="nsew")
        timer_frame.grid_columnconfigure(0, weight=1)
        timer_frame.grid_rowconfigure(0, weight=1)

        self.canvas = ctk.CTkCanvas(timer_frame, width=300, height=300, highlightthickness=0)
        self.canvas.grid(row=0, column=0, pady=20)
        
        self.mode_label = ctk.CTkLabel(timer_frame, text="Waktu Belajar", font=H2_FONT)
        self.mode_label.place(relx=0.5, rely=0.35, anchor="center")
        
        self.time_label = ctk.CTkLabel(timer_frame, text=self._format_time(self.remaining_time), font=(APP_FONT_FAMILY, 60, "bold"))
        self.time_label.place(relx=0.5, rely=0.5, anchor="center")

        controls_frame = ctk.CTkFrame(self, fg_color="transparent")
        controls_frame.grid(row=2, column=0, pady=20)
        
        self.start_pause_button = ctk.CTkButton(controls_frame, text="Mulai", font=BUTTON_FONT, width=140, height=45, command=self.toggle_timer)
        self.start_pause_button.pack(side="left", padx=10)
        
        self.reset_button = ctk.CTkButton(controls_frame, text="Reset", font=BUTTON_FONT, width=140, height=45, command=self.reset_timer, fg_color="gray50", hover_color="gray40")
        self.reset_button.pack(side="left", padx=10)

    def on_theme_change(self):
        bg_color_tuple = self.winfo_toplevel().cget("fg_color")
        current_bg_color = bg_color_tuple[0] if ctk.get_appearance_mode() == "Light" else bg_color_tuple[1]
        self.canvas.configure(bg=current_bg_color)
        self._draw_progress_bar()

    def _format_time(self, s): return f"{s//60:02d}:{s%60:02d}"

    def _draw_progress_bar(self):
        self.canvas.delete("all")
        theme_is_dark = (ctk.get_appearance_mode() == "Dark")
        theme_color_index = 1 if theme_is_dark else 0

        button_color = ctk.ThemeManager.theme["CTkButton"]["fg_color"][theme_color_index]
        border_color = ctk.ThemeManager.theme["CTkFrame"]["border_color"][theme_color_index]
        
        total_duration = int(self.study_duration_min.get() if self.current_mode == "Belajar" else self.break_duration_min.get()) * 60
        progress_color = button_color if self.current_mode == "Belajar" else SUCCESS_COLOR
        
        if total_duration == 0: total_duration = 1 
        
        self.canvas.create_oval(10, 10, 290, 290, outline=border_color, width=18)
        
        progress_percentage = self.remaining_time / total_duration
        if progress_percentage > 0:
            self.canvas.create_arc(10, 10, 290, 290, start=90, 
                                   extent=-(progress_percentage * 360), 
                                   outline=progress_color, width=18, style="arc")

    def toggle_timer(self):
        self.timer_running = not self.timer_running
        self.start_pause_button.configure(text="Jeda" if self.timer_running else "Lanjutkan")
        if self.timer_running:
            self._update_timer()
        elif self.timer_id:
            self.after_cancel(self.timer_id)
            self.timer_id = None

    def reset_timer(self, *args):
        if self.timer_id:
            self.after_cancel(self.timer_id)
            self.timer_id = None
        
        self.timer_running = False
        self.current_mode = "Belajar"
        self.mode_label.configure(text="Waktu Belajar")
        
        try:
            self.remaining_time = int(self.study_duration_min.get()) * 60
        except (ValueError, TypeError):
            self.remaining_time = 25 * 60
            self.study_duration_min.set("25")

        self.time_label.configure(text=self._format_time(self.remaining_time))
        self.start_pause_button.configure(text="Mulai", state="normal")
        self._draw_progress_bar()

    def _update_timer(self):
        if self.timer_running and self.remaining_time > 0:
            self.remaining_time -= 1
            self.time_label.configure(text=self._format_time(self.remaining_time))
            self._draw_progress_bar()
            self.timer_id = self.after(1000, self._update_timer)
        elif self.remaining_time == 0:
            self.winfo_toplevel().bell()
            self.timer_running = False
            self.current_mode = "Istirahat" if self.current_mode == "Belajar" else "Belajar"
            self.mode_label.configure(text=f"Waktu {self.current_mode}")

            duration_var = self.break_duration_min if self.current_mode == "Istirahat" else self.study_duration_min
            self.remaining_time = int(duration_var.get()) * 60

            self.time_label.configure(text=self._format_time(self.remaining_time))
            self.start_pause_button.configure(text="Mulai")
            self._draw_progress_bar()

if __name__ == "__main__":
    app = SmartStudyApp()
    app.mainloop()
