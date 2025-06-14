import customtkinter as ctk
import threading
import json
import os
import re
from tkinter import filedialog
import io

# --- Pustaka Pihak Ketiga ---
try:
    import requests
    import fitz  # PyMuPDF
    from pptx import Presentation
    from PIL import Image
    import pytesseract
except ImportError as e:
    print(f"Error: Pustaka yang dibutuhkan hilang -> {e.name}")
    print("Silakan jalankan: pip install --upgrade customtkinter requests PyMuPDF python-pptx pytesseract Pillow")
    exit()

# --- Konfigurasi Aplikasi ---
try:
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except Exception:
    print("Peringatan: Tesseract tidak ditemukan. Fitur OCR gambar tidak akan berfungsi.")

try:
    ctk.set_default_color_theme("theme.json")
except FileNotFoundError:
    print("Peringatan: file 'theme.json' tidak ditemukan.")
    ctk.set_default_color_theme("blue")

ctk.set_appearance_mode("Dark")

# Definisi Font & Warna
APP_FONT_FAMILY = "Arial"
H1_FONT = (APP_FONT_FAMILY, 48, "bold")
H2_FONT = (APP_FONT_FAMILY, 24, "bold")
H3_FONT = (APP_FONT_FAMILY, 20, "bold")
BODY_FONT = (APP_FONT_FAMILY, 14)
BUTTON_FONT = (APP_FONT_FAMILY, 16, "bold")
SMALL_FONT = (APP_FONT_FAMILY, 12)
SUCCESS_COLOR = "#00A896"
ERROR_COLOR = "#D90429"
GOLD_ACCENT_COLOR = ("#C09A3E", "#D4AF37")

# --- Model AI Gemini & Fungsi Helper ---
API_KEY = "AIzaSyAX4tK_mDxeBC4T08WS_1Iz0JSjyPSrGKo" # Otomatis diisi oleh Canvas
API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={API_KEY}"

def get_gemini_response(prompt, content_text=None):
    if not API_KEY: return "Fitur AI dinonaktifkan karena tidak ada Kunci API."
    headers = {'Content-Type': 'application/json'}
    full_prompt = f"{prompt}\n\nBerikut kontennya:\n{content_text}" if content_text else prompt
    payload = {"contents": [{"parts": [{"text": full_prompt}]}]}
    try:
        response = requests.post(API_URL, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        result = response.json()
        return result['candidates'][0]['content']['parts'][0]['text']
    except Exception as e:
        return f"Error saat menghubungi AI: {e}"

def _dynamically_resize_textbox(textbox):
    textbox.update() 
    bbox = textbox._textbox.dlineinfo("end-1c")
    if bbox:
        required_height = bbox[1] + bbox[3] + (2 * textbox._border_spacing)
        min_height = 40 
        textbox.configure(height=max(min_height, required_height))

def apply_markdown_to_textbox(textbox, text):
    textbox.configure(state="normal")
    textbox.delete("1.0", "end")
    tk_textbox = textbox._textbox
    tk_textbox.tag_config("bold", font=(APP_FONT_FAMILY, 14, "bold"))
    tk_textbox.tag_config("italic", font=(APP_FONT_FAMILY, 14, "italic"))
    segments = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
    for segment in filter(None, segments):
        if segment.startswith('**') and segment.endswith('**'):
            textbox.insert("end", segment[2:-2], "bold")
        elif segment.startswith('*') and segment.endswith('*'):
            textbox.insert("end", segment[1:-1], "italic")
        else:
            textbox.insert("end", segment)
    textbox.configure(state="disabled")
    _dynamically_resize_textbox(textbox)

def create_paragraph_bubbles(parent_frame, text):
    for widget in parent_frame.winfo_children():
        widget.destroy()
    paragraphs = text.split('\n\n')
    for para in paragraphs:
        if not para.strip(): continue
        bubble_textbox = ctk.CTkTextbox(parent_frame, font=BODY_FONT, wrap="word", corner_radius=15,
                                    fg_color=("#E5E5E5", "#3A3A3A"),
                                    border_spacing=10, activate_scrollbars=False,
                                    border_width=1, border_color=("#DCDCDC", "#4A4A4A"))
        apply_markdown_to_textbox(bubble_textbox, para.strip())
        bubble_textbox.pack(fill="x", padx=5, pady=4)

# --- Komponen-komponen UI ---
class QuizWindow(ctk.CTkToplevel):
    def __init__(self, parent, quiz_data):
        super().__init__(parent)
        self.title("Grind On")
        self.geometry("800x600")
        self.transient(parent); self.grab_set()
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1); self.grid_rowconfigure(2, weight=1)
        self.quiz_data, self.current_question_index, self.score = quiz_data, 0, 0
        self.selected_option_index = None
        self._create_widgets()
        self.display_question()
        self.bind("<Configure>", self.update_wraplength)
    def update_wraplength(self, event=None):
        new_width = self.question_frame.winfo_width() - 50 
        if new_width > 1: self.question_label.configure(wraplength=new_width)
    def _create_widgets(self):
        self.top_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.top_frame.grid(row=0, column=0, sticky="ew", padx=20, pady=10)
        self.question_number_label = ctk.CTkLabel(self.top_frame, font=H3_FONT)
        self.question_number_label.pack(side="left")
        self.score_label = ctk.CTkLabel(self.top_frame, text="Skor: 0", font=H3_FONT)
        self.score_label.pack(side="right")
        self.question_frame = ctk.CTkFrame(self, corner_radius=10)
        self.question_frame.grid(row=1, column=0, sticky="nsew", padx=20, pady=10)
        self.question_frame.grid_columnconfigure(0, weight=1); self.question_frame.grid_rowconfigure(0, weight=1)
        self.question_label = ctk.CTkLabel(self.question_frame, font=H2_FONT)
        self.question_label.grid(row=0, column=0, padx=25, pady=25, sticky="nsew")
        self.options_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.options_frame.grid(row=2, column=0, sticky="nsew", padx=20, pady=10)
        self.options_frame.grid_columnconfigure((0, 1), weight=1); self.options_frame.grid_rowconfigure((0, 1), weight=1)
        self.option_buttons = []
        for i in range(4):
            button = ctk.CTkButton(self.options_frame, font=BUTTON_FONT, height=70, corner_radius=8, command=lambda i=i: self.select_option(i))
            button.grid(row=i//2, column=i%2, padx=5, pady=5, sticky="nsew")
            self.option_buttons.append(button)
        self.bottom_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.bottom_frame.grid(row=3, column=0, sticky="ew", padx=20, pady=(10, 20))
        self.bottom_frame.grid_columnconfigure(0, weight=1)
        self.submit_button = ctk.CTkButton(self.bottom_frame, text="Submit Jawaban", font=BUTTON_FONT, height=45, command=self.submit_answer, state="disabled")
        self.submit_button.grid(row=0, column=0, sticky="ew")
        self.next_button = ctk.CTkButton(self.bottom_frame, text="Pertanyaan Berikutnya", font=BUTTON_FONT, height=45, command=self.next_question)
        self.explanation_frame = ctk.CTkFrame(self, corner_radius=10)
        self.explanation_label = ctk.CTkLabel(self.explanation_frame, text="", wraplength=750, justify="left", font=BODY_FONT)
        self.explanation_label.pack(padx=15, pady=15, fill="x")
    def display_question(self):
        self.explanation_frame.grid_forget(); self.next_button.grid_forget()
        self.submit_button.grid(row=0, column=0, sticky="ew"); self.submit_button.configure(state="disabled")
        self.selected_option_index = None
        q_data = self.quiz_data[self.current_question_index]
        self.question_number_label.configure(text=f"Soal {self.current_question_index + 1}/{len(self.quiz_data)}")
        self.question_label.configure(text=q_data.get("question", "N/A"))
        options = q_data.get("options", [])
        for i, button in enumerate(self.option_buttons):
            option_text = options[i] if i < len(options) else f"Opsi {i+1}"
            button.configure(text=option_text, state="normal", fg_color=ctk.ThemeManager.theme["CTkButton"]["fg_color"])
    def select_option(self, index):
        self.selected_option_index = index
        self.submit_button.configure(state="normal")
        for i, button in enumerate(self.option_buttons):
            button.configure(fg_color=ctk.ThemeManager.theme["CTkButton"]["hover_color"] if i == index else ctk.ThemeManager.theme["CTkButton"]["fg_color"])
    def submit_answer(self):
        q_data = self.quiz_data[self.current_question_index]
        correct_index = q_data.get("correct_answer_index")
        for button in self.option_buttons: button.configure(state="disabled")
        if self.selected_option_index == correct_index:
            self.score += 1; self.option_buttons[self.selected_option_index].configure(fg_color=SUCCESS_COLOR)
        else:
            self.option_buttons[self.selected_option_index].configure(fg_color=ERROR_COLOR)
            if correct_index is not None and 0 <= correct_index < 4:
                self.option_buttons[correct_index].configure(fg_color=SUCCESS_COLOR)
        self.score_label.configure(text=f"Skor: {self.score}")
        self.explanation_label.configure(text=f"Penjelasan:\n{q_data.get('explanation', 'N/A')}")
        self.explanation_frame.grid(row=4, column=0, padx=20, pady=10, sticky="ew")
        self.submit_button.grid_forget(); self.next_button.grid(row=0, column=0, sticky="ew")
    def next_question(self):
        self.current_question_index += 1
        if self.current_question_index < len(self.quiz_data): self.display_question()
        else: self.show_results()
    def show_results(self):
        for widget in self.winfo_children(): widget.destroy()
        self.grid_rowconfigure((0,1,2), weight=1)
        ctk.CTkLabel(self, text="Quizz Is Done!", font=(APP_FONT_FAMILY, 40, "bold")).grid(row=0, column=0)
        ctk.CTkLabel(self, text=f"Final Score: {self.score} / {len(self.quiz_data)}", font=H2_FONT).grid(row=1, column=0)
        ctk.CTkButton(self, text="Close", font=BUTTON_FONT, command=self.destroy, height=45).grid(row=2, column=0, padx=100, pady=20, sticky="ew")

class SmartStudyApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("IHTLI-Study App")
        self.geometry("1100x720"); self.minsize(900, 600)
        self.grid_columnconfigure(1, weight=1); self.grid_rowconfigure(0, weight=1)
        self._create_sidebar()
        self._create_main_frame()
        self._create_popup_frames()
        self.select_feature("pomodoro")
    def _create_sidebar(self):
        self.sidebar_frame = ctk.CTkFrame(self, width=220, corner_radius=0)
        self.sidebar_frame.grid(row=0, column=0, rowspan=4, sticky="nsew")
        self.sidebar_frame.grid_rowconfigure(6, weight=1)
        self.logo_label = ctk.CTkLabel(self.sidebar_frame, text="IHTLI-Study App", font=H2_FONT)
        self.logo_label.grid(row=0, column=0, padx=20, pady=(30, 20))
        buttons_info = {"pomodoro": "Pomodoro", "chatbot": "Ask Me AI", "summarizer": "Summarizer", "explainer": "Explainer", "quiz": "Grind On"}
        for i, (name, text) in enumerate(buttons_info.items(), start=1):
            button = ctk.CTkButton(self.sidebar_frame, text=text, font=BUTTON_FONT, anchor="w", height=45, corner_radius=8, command=lambda n=name: self.select_feature(n))
            button.grid(row=i, column=0, padx=20, pady=8, sticky="ew")
        
        # --- PERUBAHAN UI: DARI DROPDOWN KE SWITCH ---
        self.theme_switch = ctk.CTkSwitch(self.sidebar_frame, text="Theme", font=BODY_FONT, command=self.toggle_theme)
        self.theme_switch.grid(row=8, column=0, padx=20, pady=(20, 20), sticky="w")
        self.theme_switch.select() # Mulai dalam mode gelap

    def toggle_theme(self):
        # Fungsi baru untuk menangani logika switch
        if self.theme_switch.get() == 1:
            ctk.set_appearance_mode("Dark")
        else:
            ctk.set_appearance_mode("Light")
        
        # Update custom colored widgets
        if "pomodoro" in self.popup_frames and self.popup_frames["pomodoro"].winfo_exists():
             self.popup_frames["pomodoro"].on_theme_change()

    def _create_main_frame(self):
        self.main_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.main_frame.grid(row=0, column=1, padx=20, pady=20, sticky="nsew")
        self.main_frame.grid_columnconfigure(0, weight=1); self.main_frame.grid_rowconfigure(1, weight=1)
        self.feature_title_frame = ctk.CTkFrame(self.main_frame, corner_radius=10, height=60)
        self.feature_title_frame.grid(row=0, column=0, pady=(0, 10), sticky="ew")
        self.feature_title_label = ctk.CTkLabel(self.feature_title_frame, text="", font=H2_FONT)
        self.feature_title_label.pack(side="left", padx=20)
    def _create_popup_frames(self):
        self.popup_frames = {
            "pomodoro": PomodoroFrame(self.main_frame), "chatbot": ChatbotFrame(self.main_frame),
            "summarizer": DocumentFeatureFrame(self.main_frame, "summarize"),
            "explainer": DocumentFeatureFrame(self.main_frame, "explain"),
            "quiz": DocumentFeatureFrame(self.main_frame, "quiz")}
    def select_feature(self, name):
        titles = {"pomodoro": "Pomodoro", "chatbot": "Ask Me AI", "summarizer": "Summarizer", "explainer": "Explainer", "quiz": "Grind On"}
        self.feature_title_label.configure(text=titles.get(name, "Fitur"))
        for frame_name, frame in self.popup_frames.items():
            if frame_name == name:
                self._animate_slide_and_fade_in(frame)
            else:
                frame.grid_forget()
    def _animate_slide_and_fade_in(self, widget, steps=20, duration_ms=300):
        widget.grid(row=1, column=0, sticky="nsew", pady=10)
        widget.update_idletasks()
        start_y, end_y = widget.winfo_y() + 20, widget.winfo_y()
        try:
            widget.attributes("-alpha", 0); widget.lift()
            def animation_step(step):
                if step > steps:
                    widget.place_forget(); widget.grid(row=1, column=0, sticky="nsew", pady=0)
                    widget.attributes("-alpha", 1); return
                progress = step / steps
                eased_progress = 1 - pow(1 - progress, 3)
                current_y = start_y - (start_y - end_y) * eased_progress
                widget.place(y=current_y, relx=0, rely=0, relwidth=1.0)
                widget.attributes("-alpha", eased_progress)
                self.after(duration_ms // steps, lambda: animation_step(step + 1))
            animation_step(0)
        except Exception: widget.lift()

class ChatbotFrame(ctk.CTkFrame):
    def __init__(self, parent):
        super().__init__(parent, fg_color="transparent")
        self.grid_columnconfigure(0, weight=1); self.grid_rowconfigure(0, weight=1)
        self.chat_history_frame = ctk.CTkScrollableFrame(self, corner_radius=10)
        self.chat_history_frame.grid(row=0, column=0, sticky="nsew", pady=(0,10))
        self.chat_history_frame.grid_columnconfigure(0, weight=1)
        self.input_frame = ctk.CTkFrame(self, corner_radius=10)
        self.input_frame.grid(row=1, column=0, sticky="ew")
        self.input_frame.grid_columnconfigure(0, weight=1)
        self.user_input = ctk.CTkEntry(self.input_frame, placeholder_text="Ketik pesan Anda...", font=BODY_FONT, border_width=0, fg_color="transparent")
        self.user_input.grid(row=0, column=0, padx=15, pady=5, sticky="ew")
        self.user_input.bind("<Return>", lambda e: self.send_message())
        self.send_button = ctk.CTkButton(self.input_frame, text="Kirim", font=BUTTON_FONT, width=80, height=35, command=self.send_message)
        self.send_button.grid(row=0, column=1, padx=(0, 5), pady=5)
        self._add_message_bubble("bot", "Hello! Im an AI assistant integgrated with Gemini 2.0, how can I help you?")
    def _add_message_bubble(self, sender, message):
        is_user = sender == "user"
        if is_user:
            self._create_user_bubble(message)
        else:
            create_paragraph_bubbles(self.chat_history_frame, message)
        self.after(100, self.chat_history_frame._parent_canvas.yview_moveto, 1.0)
    def _create_user_bubble(self, message):
        msg_frame = ctk.CTkFrame(self.chat_history_frame, fg_color="transparent")
        msg_frame.pack(fill="x", padx=5, pady=(5, 1))
        msg_bubble = ctk.CTkTextbox(msg_frame, font=BODY_FONT, wrap="word", corner_radius=15,
                                    fg_color=GOLD_ACCENT_COLOR, text_color=("#1A1A1A", "#1A1A1A"),
                                    border_spacing=10, activate_scrollbars=False)
        msg_bubble.insert("end", message)
        msg_bubble.configure(state="disabled")
        _dynamically_resize_textbox(msg_bubble)
        msg_bubble.pack(anchor="e", padx=(80, 5), pady=5)
    def send_message(self):
        user_message = self.user_input.get()
        if not user_message.strip(): return
        self._add_message_bubble("user", user_message)
        self.user_input.delete(0, "end")
        threading.Thread(target=self._get_ai_response, args=(user_message,), daemon=True).start()
    def _get_ai_response(self, message):
        prompt = "Anda adalah asisten yang ramah dan cerdas. Jawab pertanyaan ini dengan jelas. Gunakan markdown **bold** dan *italic* untuk penekanan. Jika jawaban Anda panjang, pastikan untuk memisahkannya menjadi beberapa paragraf dengan memberikan dua baris baru (dua kali enter) di antara paragraf."
        response = get_gemini_response(prompt, message)
        self.after(0, lambda: self._add_message_bubble("bot", response))

class DocumentFeatureFrame(ctk.CTkFrame):
    def __init__(self, parent, feature_type):
        super().__init__(parent, fg_color="transparent")
        self.grid_columnconfigure(0, weight=1); self.grid_rowconfigure(1, weight=1)
        self.parent, self.feature_type, self.file_content = parent, feature_type, None
        self._create_widgets()
    def _create_widgets(self):
        top_frame = ctk.CTkFrame(self, corner_radius=10)
        top_frame.grid(row=0, column=0, pady=(0, 10), sticky="ew")
        top_frame.grid_columnconfigure(1, weight=1)
        self.upload_button = ctk.CTkButton(top_frame, text="Upload File", font=BUTTON_FONT, height=40, command=self.upload_file)
        self.upload_button.grid(row=0, column=0, padx=10, pady=10)
        self.file_label = ctk.CTkLabel(top_frame, text="Import File", text_color="gray", font=SMALL_FONT, anchor="w")
        self.file_label.grid(row=0, column=1, padx=10, sticky="ew")
        button_text = "Mulai Kuis" if self.feature_type == 'quiz' else "Proses Dokumen"
        self.run_button = ctk.CTkButton(top_frame, text=button_text, font=BUTTON_FONT, height=40, command=self.run_feature, state="disabled")
        self.run_button.grid(row=0, column=2, padx=10, pady=10)
        self.output_scroll_frame = ctk.CTkScrollableFrame(self, corner_radius=10)
        self.output_scroll_frame.grid(row=1, column=0, sticky="nsew")
        self.output_scroll_frame.grid_columnconfigure(0, weight=1)
        self.initial_label = ctk.CTkLabel(self.output_scroll_frame, text="The document is need to be uploaded to start AI analysis.", font=BODY_FONT, text_color="gray")
        self.initial_label.pack(expand=True, padx=20, pady=20)
    def _set_output(self, text):
        self.initial_label.pack_forget()
        create_paragraph_bubbles(self.output_scroll_frame, text)
    def upload_file(self):
        filetypes = (("Dokumen & Gambar", "*.pdf *.pptx *.png *.jpg *.jpeg *.bmp *.tiff"),)
        filepath = filedialog.askopenfilename(filetypes=filetypes)
        if not filepath: return
        self.file_label.configure(text=os.path.basename(filepath), text_color=ctk.ThemeManager.theme["CTkLabel"]["text_color"])
        self._set_output(f"Menganalisis file '{os.path.basename(filepath)}'...")
        self.run_button.configure(state="disabled")
        threading.Thread(target=self._extract_text, args=(filepath,), daemon=True).start()
    def _extract_text(self, file_path):
        try:
            ext = os.path.splitext(file_path)[1].lower()
            text = ""
            if ext == '.pdf':
                with fitz.open(file_path) as doc:
                    if doc.is_encrypted: raise ValueError("PDF ini dilindungi kata sandi.")
                    for page in doc:
                        text += page.get_text("text", sort=True) + "\n"
                        for img_index, img in enumerate(page.get_images(full=True)):
                            try:
                                ocr_text = pytesseract.image_to_string(Image.open(io.BytesIO(doc.extract_image(img[0])["image"])))
                                if ocr_text.strip(): text += f"\n--- Teks dari Gambar ---\n{ocr_text}\n"
                            except Exception: pass
            elif ext == '.pptx':
                text = "\n".join(s.text for slide in Presentation(file_path).slides for s in slide.shapes if hasattr(s, "text"))
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                text = pytesseract.image_to_string(Image.open(file_path))
            else: raise ValueError(f"Format file tidak didukung: {ext}")
            cleaned_text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
            if not cleaned_text or len(cleaned_text) < 20: raise ValueError("Tidak dapat menemukan teks signifikan dalam file.")
            self.file_content = cleaned_text
            self.after(0, self._on_extraction_complete, os.path.basename(file_path))
        except Exception as e: self.after(0, self._set_output, f"**Gagal memproses file:**\n\n{e}")
    def _on_extraction_complete(self, filename):
        action = "memulai kuis" if self.feature_type == 'quiz' else "memprosesnya"
        self._set_output(f"**Sukses!** File '{filename}' telah dianalisis.\n\nAnda sekarang bisa *{action}*.")
        self.run_button.configure(state="normal")
    def run_feature(self):
        if not self.file_content: self._set_output("**Error:** Tidak ada konten file."); return
        self._set_output("Please Wait.")
        self.run_button.configure(state="disabled")
        threading.Thread(target=self._get_ai_result, daemon=True).start()
    def _get_ai_result(self):
        prompts = {
            "summarize": "Anda adalah seorang analis ahli. Buat rangkuman yang **sangat mendalam dan terperinci** dari teks berikut. Bagi menjadi beberapa bagian dengan **sub-judul tebal**. Jelaskan setiap poin utama dengan beberapa kalimat, jangan hanya satu baris. Pastikan semua informasi krusial tercakup dan pisahkan setiap bagian menjadi paragraf yang jelas.",
            "explain": "Anda adalah seorang dosen ahli. Identifikasi **semua konsep penting** dalam teks berikut. Untuk setiap konsep, berikan **penjelasan yang komprehensif**, termasuk definisi, contoh, dan analogi. Pastikan jawaban Anda terstruktur dalam beberapa paragraf.",
            "quiz": "Anda adalah pembuat kuis. Buat 10 pertanyaan dari teks ini. Respons HARUS HANYA dalam format JSON array. Setiap objek berisi: 'question' (string), 'options' (array 4 string), 'correct_answer_index' (int 0-3), dan 'explanation' (string)."
        }
        prompt = prompts.get(self.feature_type)
        response_text = get_gemini_response(prompt, self.file_content)
        if self.feature_type == 'quiz':
            try:
                if "```json" in response_text: response_text = response_text.split("```json")[1].split("```")[0]
                quiz_data = json.loads(response_text)
                if not isinstance(quiz_data, list) or len(quiz_data) < 1: raise ValueError("Respons JSON tidak valid.")
                self.after(0, lambda: QuizWindow(self.parent, quiz_data))
                self.after(0, self._set_output, "**Sukses!** Kuis telah dimulai di jendela baru.")
            except Exception as e:
                self.after(0, self._set_output, f"**Gagal membuat kuis.**\n\nAI memberikan respons tidak valid.\nDetail: {e}")
            finally:
                self.after(0, self.run_button.configure, {"state": "normal"})
        else:
            self.after(0, lambda: (self._set_output(response_text), self.run_button.configure(state="normal")))

class PomodoroFrame(ctk.CTkFrame):
    def __init__(self, parent):
        super().__init__(parent, fg_color="transparent")
        self.grid_columnconfigure(0, weight=1); self.grid_rowconfigure(1, weight=1)
        self.timer_running, self.timer_id, self.current_mode = False, None, "Belajar"
        self.study_duration_min = ctk.StringVar(value="25"); self.break_duration_min = ctk.StringVar(value="5")
        self.remaining_time = int(self.study_duration_min.get()) * 60
        self._create_widgets()
        self.on_theme_change()
    def _create_widgets(self):
        settings_frame = ctk.CTkFrame(self, fg_color="transparent")
        settings_frame.grid(row=0, column=0, pady=(0, 20))
        ctk.CTkLabel(settings_frame, text="Study Time:", font=BODY_FONT).pack(side="left", padx=(0, 5))
        ctk.CTkOptionMenu(settings_frame, variable=self.study_duration_min, font=BODY_FONT, values=[str(i) for i in [5, 15, 25, 30, 45, 60]], command=self.reset_timer).pack(side="left", padx=(0, 20))
        ctk.CTkLabel(settings_frame, text="Break Time:", font=BODY_FONT).pack(side="left", padx=(0, 5))
        ctk.CTkOptionMenu(settings_frame, variable=self.break_duration_min, font=BODY_FONT, values=[str(i) for i in [1, 5, 10, 15, 20]], command=self.reset_timer).pack(side="left")
        timer_frame = ctk.CTkFrame(self, fg_color="transparent")
        timer_frame.grid(row=1, column=0, sticky="nsew")
        timer_frame.grid_columnconfigure(0, weight=1); timer_frame.grid_rowconfigure(0, weight=1)
        self.canvas = ctk.CTkCanvas(timer_frame, width=300, height=300, highlightthickness=0)
        self.canvas.grid(row=0, column=0, pady=20)
        self.mode_label = ctk.CTkLabel(timer_frame, text="Study Time", font=H2_FONT)
        self.mode_label.place(relx=0.5, rely=0.35, anchor="center")
        self.time_label = ctk.CTkLabel(timer_frame, text=self._format_time(self.remaining_time), font=(APP_FONT_FAMILY, 60, "bold"))
        self.time_label.place(relx=0.5, rely=0.5, anchor="center")
        controls_frame = ctk.CTkFrame(self, fg_color="transparent")
        controls_frame.grid(row=2, column=0, pady=20)
        self.start_pause_button = ctk.CTkButton(controls_frame, text="Start", font=BUTTON_FONT, width=140, height=45, command=self.toggle_timer)
        self.start_pause_button.pack(side="left", padx=10)
        self.reset_button = ctk.CTkButton(controls_frame, text="Reset", font=BUTTON_FONT, width=140, height=45, command=self.reset_timer, fg_color="gray50", hover_color="gray40")
        self.reset_button.pack(side="left", padx=10)
    def on_theme_change(self):
        # --- PERBAIKAN WARNA ---
        # Mengambil warna dari root window ("CTk") bukan "CTkFrame"
        color_tuple = ctk.ThemeManager.theme["CTk"]["fg_color"]
        self.canvas.configure(bg=color_tuple[0] if ctk.get_appearance_mode() == "Light" else color_tuple[1])
        self._draw_progress_bar()
    def _format_time(self, s): return f"{s//60:02d}:{s%60:02d}"
    def _draw_progress_bar(self):
        self.canvas.delete("all")
        theme_color_index = 1 if ctk.get_appearance_mode() == "Dark" else 0
        button_colors = ctk.ThemeManager.theme["CTkButton"]["fg_color"]
        study_progress_color = button_colors[theme_color_index]
        border_color_tuple = ctk.ThemeManager.theme["CTkFrame"]["border_color"]
        border_color = border_color_tuple[theme_color_index]
        total_duration = int(self.study_duration_min.get() if self.current_mode == "Belajar" else self.break_duration_min.get()) * 60
        progress_color = study_progress_color if self.current_mode == "Belajar" else SUCCESS_COLOR
        if total_duration == 0: total_duration = 1
        self.canvas.create_oval(10, 10, 290, 290, outline=border_color, width=18)
        if (self.remaining_time / total_duration) > 0:
            self.canvas.create_arc(10, 10, 290, 290, start=90, extent=-(self.remaining_time / total_duration * 360), outline=progress_color, width=18, style="arc")
    def toggle_timer(self):
        self.timer_running = not self.timer_running
        self.start_pause_button.configure(text="Jeda" if self.timer_running else "Lanjutkan")
        if self.timer_running: self._update_timer()
        elif self.timer_id: self.after_cancel(self.timer_id)
    def reset_timer(self, *args):
        if self.timer_id: self.after_cancel(self.timer_id)
        self.timer_running = False
        self.current_mode = "Belajar"
        self.mode_label.configure(text="Waktu Belajar")
        self.remaining_time = int(self.study_duration_min.get()) * 60
        self.time_label.configure(text=self._format_time(self.remaining_time))
        self.start_pause_button.configure(text="Mulai", state="normal"); self._draw_progress_bar()
    def _update_timer(self):
        if self.timer_running and self.remaining_time > 0:
            self.remaining_time -= 1; self.time_label.configure(text=self._format_time(self.remaining_time)); self._draw_progress_bar()
            self.timer_id = self.after(1000, self._update_timer)
        elif self.remaining_time == 0:
            self.timer_running = False
            self.current_mode = "Istirahat" if self.current_mode == "Belajar" else "Belajar"
            self.mode_label.configure(text=f"Waktu {self.current_mode}")
            self.remaining_time = int(self.break_duration_min.get() if self.current_mode == "Istirahat" else self.study_duration_min.get()) * 60
            self.time_label.configure(text=self._format_time(self.remaining_time))
            self.start_pause_button.configure(text="Mulai"); self._draw_progress_bar()

if __name__ == "__main__":
    app = SmartStudyApp()
    app.mainloop()